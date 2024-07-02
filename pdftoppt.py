import streamlit as st
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
import io
from PIL import Image

def pdf_to_ppt(pdf_bytes):
    # Create a presentation object
    prs = Presentation()
    
    # Open the PDF
    pdf_document = fitz.open(stream=pdf_bytes, filetype="pdf")
    
    for page_num in range(len(pdf_document)):
        page = pdf_document.load_page(page_num)
        
        # Create a slide
        slide_layout = prs.slide_layouts[5]  # Use a blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Extract text
        text = page.get_text("text")
        
        # Add text to the slide
        left = Inches(1)
        top = Inches(1)
        width = Inches(8.5)
        height = Inches(5.5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = text
        
        # Extract images
        images = page.get_images(full=True)
        for img_index, img in enumerate(images):
            xref = img[0]
            base_image = pdf_document.extract_image(xref)
            image_bytes = base_image["image"]
            image_ext = base_image["ext"]
            image = Image.open(io.BytesIO(image_bytes))
            
            # Save image to a temporary file in memory
            image_stream = io.BytesIO()
            image.save(image_stream, format=image_ext)
            image_stream.seek(0)
            
            # Add image to slide
            left = Inches(1)
            top = Inches(1 + img_index * 3)  # Adjust the position of each image
            slide.shapes.add_picture(image_stream, left, top, width=Inches(6))
    
    # Save the presentation to a BytesIO object
    ppt_bytes = io.BytesIO()
    prs.save(ppt_bytes)
    ppt_bytes.seek(0)
    return ppt_bytes

st.title("PDF to PowerPoint Converter")

uploaded_file = st.file_uploader("Upload a PDF file", type=["pdf"])

if uploaded_file is not None:
    ppt_bytes = pdf_to_ppt(uploaded_file.getvalue())
    st.success("Conversion successful! Download your PowerPoint file below.")
    st.download_button(
        label="Download PowerPoint",
        data=ppt_bytes,
        file_name="converted_presentation.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
